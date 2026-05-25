#!/usr/bin/env python3
"""
create_presentation.py
Generate docs/Apresentacao_Pipeline_GFET.pptx (9 slides, Portuguese).
Run: python create_presentation.py
Requires: pip install python-pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ---------------------------------------------------------------------------
# CONSTANTS
# ---------------------------------------------------------------------------
W = 13.333   # slide width  (inches, 16:9)
H = 7.5      # slide height (inches)

NAVY       = RGBColor(0x1F, 0x38, 0x64)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLUE_ACC   = RGBColor(0x2E, 0x75, 0xB6)
BLUE_LIGHT = RGBColor(0xDD, 0xEA, 0xF5)
GREY_BG    = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GREY  = RGBColor(0x26, 0x26, 0x26)
MID_GREY   = RGBColor(0xD9, 0xD9, 0xD9)
BLACK      = RGBColor(0x00, 0x00, 0x00)
YELLOW_BG  = RGBColor(0xFF, 0xF2, 0xCC)
ORANGE_BG  = RGBColor(0xFC, 0xE4, 0xD6)
GREEN_MED  = RGBColor(0xC6, 0xE0, 0xB4)

ROW_HEADER = BLUE_ACC
ROW_GOOD   = RGBColor(0xE2, 0xEF, 0xDA)
ROW_MOD    = RGBColor(0xFF, 0xF2, 0xCC)
ROW_LOW    = RGBColor(0xFF, 0xE0, 0xDD)
ROW_ALT    = RGBColor(0xEA, 0xF2, 0xFB)

# ---------------------------------------------------------------------------
# HELPERS
# ---------------------------------------------------------------------------

def I(n):  return Inches(n)
def P(n):  return Pt(n)


def new_prs():
    prs = Presentation()
    prs.slide_width  = I(W)
    prs.slide_height = I(H)
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def set_bg(slide, color=GREY_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, bg=None, line_color=None):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, I(left), I(top), I(width), I(height)
    )
    if bg:
        shp.fill.solid()
        shp.fill.fore_color.rgb = bg
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
    else:
        shp.line.fill.background()
    return shp


def add_text_box(slide, left, top, width, height, text, size=13,
                 bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT,
                 wrap=True, italic=False):
    txb = slide.shapes.add_textbox(I(left), I(top), I(width), I(height))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = P(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_title_bar(slide, title, subtitle=None):
    bar = add_rect(slide, 0, 0, W, 1.1, bg=NAVY)
    tf  = bar.text_frame
    tf.word_wrap = False
    tf.margin_left   = I(0.3)
    tf.margin_top    = I(0.13)
    tf.margin_bottom = I(0.05)
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size  = P(26)
    run.font.bold  = True
    run.font.color.rgb = WHITE
    if subtitle:
        p2  = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2  = p2.add_run()
        r2.text = subtitle
        r2.font.size  = P(13)
        r2.font.color.rgb = RGBColor(0xBD, 0xD7, 0xEE)
    return bar


def add_step_box(slide, left, top, width, height, num, title, desc, bg):
    shp = add_rect(slide, left, top, width, height, bg=bg, line_color=MID_GREY)
    tf  = shp.text_frame
    tf.word_wrap = True
    tf.margin_left   = I(0.12)
    tf.margin_top    = I(0.08)
    tf.margin_right  = I(0.08)
    tf.margin_bottom = I(0.05)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"{num}. {title}"
    r.font.size  = P(13)
    r.font.bold  = True
    r.font.color.rgb = DARK_GREY
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = desc
    r2.font.size  = P(10)
    r2.font.color.rgb = RGBColor(0x40, 0x40, 0x40)
    return shp


def add_bullet_text(slide, left, top, width, height, lines, size=13,
                    bold_first=False, leading_size=None):
    txb = slide.shapes.add_textbox(I(left), I(top), I(width), I(height))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = P(2)
        r = p.add_run()
        r.text = line
        r.font.size  = P(leading_size if (i == 0 and leading_size) else size)
        r.font.bold  = bold_first and i == 0
        r.font.color.rgb = DARK_GREY
    return txb


def make_table(slide, left, top, width, height, headers, rows,
               header_bg=ROW_HEADER, header_text=WHITE,
               row_colors=None, col_widths=None, header_size=10, row_size=9):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shp = slide.shapes.add_table(n_rows, n_cols, I(left), I(top), I(width), I(height))
    tbl = tbl_shp.table

    # Column widths
    if col_widths:
        total = sum(col_widths)
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = int(I(width) * cw / total)

    # Header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        _set_cell_text(cell, hdr, size=header_size, bold=True,
                       color=header_text, align=PP_ALIGN.CENTER)

    # Data rows
    for ri, row in enumerate(rows):
        bg = None
        if row_colors:
            bg = row_colors[ri]
        elif ri % 2 == 1:
            bg = ROW_ALT
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if bg:
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
            _set_cell_text(cell, str(val), size=row_size,
                           align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    return tbl_shp


def _set_cell_text(cell, text, size=9, bold=False,
                   color=DARK_GREY, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.margin_top    = P(1.5)
    tf.margin_bottom = P(1.5)
    tf.margin_left   = P(3)
    tf.margin_right  = P(3)
    p = tf.paragraphs[0]
    p.alignment = align
    # clear default run if present
    for run in p.runs:
        run.text = ""
    r = p.add_run()
    r.text = text
    r.font.size  = P(size)
    r.font.bold  = bold
    r.font.color.rgb = color


def note_box(slide, left, top, width, height, text, size=11):
    shp = add_rect(slide, left, top, width, height,
                   bg=YELLOW_BG, line_color=RGBColor(0xC9, 0xA8, 0x00))
    tf  = shp.text_frame
    tf.word_wrap = True
    tf.margin_left  = I(0.1)
    tf.margin_top   = I(0.07)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size  = P(size)
    r.font.color.rgb = RGBColor(0x40, 0x30, 0x00)
    return shp


# ---------------------------------------------------------------------------
# SLIDE 1 — TITLE
# ---------------------------------------------------------------------------
def slide_title(prs):
    slide = blank_slide(prs)
    set_bg(slide, NAVY)

    # Big white title
    add_text_box(slide, 0.5, 1.4, 12.3, 2.0,
                 "Desenho Computacional de\nProbes de DNA para\nBiossensores GFET",
                 size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle bar
    bar = add_rect(slide, 0, 4.05, W, 0.95, bg=BLUE_ACC)
    add_text_box(slide, 0, 4.12, W, 0.8,
                 "Pipeline in silico para deteção de 6 patógenos bacterianos",
                 size=17, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    # Author info
    add_text_box(slide, 0, 5.3, W, 1.5,
                 "Romeu Fernandes  ·  PG45861  ·  Universidade do Minho  ·  Maio 2026",
                 size=14, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.CENTER)

    # Bottom accent line
    add_rect(slide, 0, H - 0.15, W, 0.15, bg=RGBColor(0xFF, 0xC0, 0x00))


# ---------------------------------------------------------------------------
# SLIDE 2 — ALVOS BIOLÓGICOS
# ---------------------------------------------------------------------------
def slide_targets(prs):
    slide = blank_slide(prs)
    set_bg(slide, GREY_BG)
    add_title_bar(slide, "Alvos Biológicos", "Genes-alvo selecionados para deteção por GFET")

    # Left column: GFET context
    add_text_box(slide, 0.3, 1.25, 5.5, 0.4,
                 "Biossensores GFET — Contexto", size=14, bold=True, color=NAVY)
    bullets_left = [
        "• GFET (Graphene Field-Effect Transistor): biossensor de elevada\n"
        "  sensibilidade baseado em grafeno",
        "• Probe de ssDNA imobilizada na superfície de grafeno capta o\n"
        "  alvo via hibridação → mudança de carga → sinal elétrico",
        "• Objetivo: desenhar probes com Tm, GC% e estrutura secundária\n"
        "  ótimos para 6 genes bacterianos",
    ]
    add_bullet_text(slide, 0.3, 1.65, 5.5, 2.8, bullets_left, size=12)

    # Right: targets table
    add_text_box(slide, 6.2, 1.25, 6.8, 0.4,
                 "Genes Alvo", size=14, bold=True, color=NAVY)

    headers = ["Gene", "Organismo", "Grupo"]
    rows = [
        ("nuc",   "Staphylococcus aureus",     "A"),
        ("rmpM",  "Neisseria meningitidis",     "A"),
        ("lytA",  "Streptococcus pneumoniae",   "B"),
        ("oprL",  "Pseudomonas aeruginosa",     "B"),
        ("algD",  "Pseudomonas aeruginosa",     "B"),
        ("frdB",  "Haemophilus influenzae",     "B"),
    ]
    row_colors = [ROW_ALT if i % 2 == 1 else None for i in range(6)]
    # mark algD row
    row_colors[4] = RGBColor(0xFF, 0xE0, 0xDD)

    make_table(slide, 6.2, 1.65, 6.8, 3.2,
               headers, rows,
               col_widths=[1.2, 4.2, 1.4],
               row_colors=row_colors,
               header_size=11, row_size=11)

    note_box(slide, 6.2, 5.0, 6.8, 0.55,
             "algD: gene problemático (ver Slide 4) — apenas 2 sequências no NCBI",
             size=10)

    # Source note
    add_text_box(slide, 0.3, 6.5, 12.5, 0.4,
                 "Grupo A: patógenos Gram-positivos / Gram-negativos sistémicos  |  "
                 "Grupo B: patógenos respiratórios / Gram-negativos",
                 size=10, color=RGBColor(0x60, 0x60, 0x60), italic=True)


# ---------------------------------------------------------------------------
# SLIDE 3 — VISÃO GERAL DA PIPELINE
# ---------------------------------------------------------------------------
def slide_pipeline_overview(prs):
    slide = blank_slide(prs)
    set_bg(slide, GREY_BG)
    add_title_bar(slide, "Visão Geral da Pipeline",
                  "Sequência de execução — do NCBI à previsão 3D")

    # Step box definitions
    steps = [
        (1, "NCBI Entrez",           "Recuperação de sequências\n≤20 seq/gene, formato FASTA",
         RGBColor(0xDD, 0xEA, 0xF5)),
        (2, "MAFFT v7.526",          "Alinhamento múltiplo\nde sequências (MSA)",
         RGBColor(0xC5, 0xDA, 0xF0)),
        (3, "Janelas Conservadas",   "18–28 nt · conservação ≥ 85%\ngaps ≤ 20%  →  candidatas",
         RGBColor(0xAD, 0xCB, 0xEB)),
        (4, "Triagem Termoquímica",  "Tm, GC%, hairpin ΔG,\nhomodimer ΔG  (primer3-py)",
         ROW_GOOD),
        (5, "seqfold (MFE)",         "Estrutura secundária ssDNA\nΔG MFE > −6.0 kcal/mol",
         GREEN_MED),
        (6, "Boltz-2 (3D)",          "Previsão de estrutura 3D\nconfidence · pTM · pLDDT",
         RGBColor(0xFC, 0xE4, 0xD6)),
    ]

    box_w, box_h = 3.6, 1.4
    row1_y = 1.8
    row2_y = 4.55
    arrow_y_r1 = row1_y + box_h / 2 - 0.18
    arrow_y_r2 = row2_y + box_h / 2 - 0.18
    col_x = [0.25, 4.65, 9.05]

    for i, (num, title, desc, bg) in enumerate(steps):
        row = 0 if i < 3 else 1
        col = i if i < 3 else i - 3
        y   = row1_y if row == 0 else row2_y
        x   = col_x[col]
        add_step_box(slide, x, y, box_w, box_h, num, title, desc, bg)

    # Horizontal arrows row 1
    for ci in range(2):
        ax = col_x[ci] + box_w + 0.04
        add_text_box(slide, ax, arrow_y_r1, 0.52, 0.45, "→",
                     size=22, bold=True, color=BLUE_ACC, align=PP_ALIGN.CENTER)

    # Horizontal arrows row 2
    for ci in range(2):
        ax = col_x[ci] + box_w + 0.04
        add_text_box(slide, ax, arrow_y_r2, 0.52, 0.45, "→",
                     size=22, bold=True, color=BLUE_ACC, align=PP_ALIGN.CENTER)

    # Down arrow between row 1 and row 2 (right side)
    down_x = col_x[2] + box_w + 0.08
    down_y = row1_y + box_h - 0.1
    add_text_box(slide, down_x, down_y, 0.7, 1.35, "↓",
                 size=26, bold=True, color=BLUE_ACC, align=PP_ALIGN.CENTER)

    # Output label
    note_box(slide, 0.25, 6.2, 12.8, 0.55,
             "Output: FINAL_PROBES_ALL.csv (4 263 probes geradas pela pipeline + "
             "75 probes de literatura) · Top 25 → Boltz-2 (Google Colab GPU T4)",
             size=11)


# ---------------------------------------------------------------------------
# SLIDE 4 — PASSOS 1–3: RECUPERAÇÃO, ALINHAMENTO E CANDIDATOS
# ---------------------------------------------------------------------------
def slide_steps123(prs):
    slide = blank_slide(prs)
    set_bg(slide, GREY_BG)
    add_title_bar(slide, "Passos 1–3: Recuperação, Alinhamento e Candidatos",
                  "NCBI Entrez → MAFFT → Janelas deslizantes conservadas")

    # Left panel: methodology
    add_text_box(slide, 0.3, 1.2, 5.8, 0.35, "Metodologia", size=14, bold=True, color=NAVY)
    method_lines = [
        "Passo 1 — NCBI Entrez",
        "  • Query: gene[Gene Name] + organism[Organism]",
        "  • Máximo 20 sequências por gene (RefSeq preferido)",
        "  • Filtro de comprimento: 100–5000 bp",
        "",
        "Passo 2 — MAFFT v7.526",
        "  • Alinhamento múltiplo progressivo (FFT-NS-2)",
        "  • Modo single-sequence para genes com 1 seq",
        "",
        "Passo 3 — Janelas Deslizantes Conservadas",
        "  • Comprimento: 18–28 nt  (Wetmur 1991)",
        "  • Passo: 3 nt",
        "  • Conservação por posição ≥ 85%",
        "  • Conteúdo de gaps ≤ 20%",
    ]
    add_bullet_text(slide, 0.3, 1.6, 5.8, 4.8, method_lines, size=11)

    # Right panel: per-gene stats table
    add_text_box(slide, 6.5, 1.2, 6.5, 0.35,
                 "Sequências recuperadas e candidatos gerados",
                 size=13, bold=True, color=NAVY)

    headers = ["Gene", "Organismo (abrev.)", "Seq. NCBI", "Comprimento (bp)", "Candidatos"]
    rows = [
        ("nuc",   "S. aureus",    "20",  "200–3 000",  "529"),
        ("rmpM",  "N. meningitidis", "20", "100–1 200", "76"),
        ("lytA",  "S. pneumoniae","20",  "700–1 300",  "925"),
        ("oprL",  "P. aeruginosa","20",  "300–2 000",  "1 015"),
        ("algD",  "P. aeruginosa", "2",  "500–2 500",  "2"),
        ("frdB",  "H. influenzae","20",  "200–2 000",  "1 716"),
        ("TOTAL", "—",            "102", "—",           "4 263"),
    ]
    total_bg = RGBColor(0xBD, 0xD7, 0xEE)
    row_colors = [ROW_ALT if i % 2 == 1 else None for i in range(7)]
    row_colors[4] = ROW_LOW   # algD
    row_colors[6] = total_bg  # total

    make_table(slide, 6.5, 1.6, 6.5, 3.7,
               headers, rows,
               col_widths=[0.9, 2.4, 1.1, 1.7, 1.3],
               row_colors=row_colors,
               header_size=10, row_size=10)

    note_box(slide, 6.5, 5.45, 6.5, 0.75,
             "algD: apenas 2 sequências encontradas no NCBI (P. aeruginosa).\n"
             "Resultado: 0 candidatos → gene excluído da análise estrutural.",
             size=10)


# ---------------------------------------------------------------------------
# SLIDE 5 — PASSO 4: TRIAGEM TERMOQUÍMICA
# ---------------------------------------------------------------------------
def slide_thermo(prs):
    slide = blank_slide(prs)
    set_bg(slide, GREY_BG)
    add_title_bar(slide, "Passo 4: Triagem Termoquímica",
                  "Cálculo de Tm, GC%, hairpin ΔG e homodimer ΔG via primer3-py")

    # Left: criteria table
    add_text_box(slide, 0.3, 1.2, 5.8, 0.35, "Critérios de Filtragem", size=14, bold=True, color=NAVY)

    h_crit = ["Parâmetro", "Critério (default)", "Referência"]
    rows_crit = [
        ("Tm",            "53–72 °C",          "SantaLucia & Hicks 2004"),
        ("GC%",           "40–65%",            "IDT OligoAnalyzer"),
        ("Hairpin ΔG",    "> −2 kcal/mol",     "primer3-py"),
        ("Homodimer ΔG",  "> −5 kcal/mol",     "primer3-py"),
        ("Comprimento",   "18–28 nt",          "Wetmur 1991"),
    ]
    make_table(slide, 0.3, 1.6, 5.8, 2.6,
               h_crit, rows_crit,
               col_widths=[2.0, 1.8, 2.0],
               header_size=11, row_size=11)

    add_text_box(slide, 0.3, 4.3, 5.8, 0.35,
                 "Ajustes gene-específicos", size=12, bold=True, color=NAVY)
    add_bullet_text(slide, 0.3, 4.65, 5.8, 1.8,
                    ["• nuc, frdB (AT-rich): Tm mín. 52 °C · GC 38–60%",
                     "• oprL, algD (GC-rich): GC máx. 70%",
                     "• rmpM: conservação mín. 80% (variabilidade elevada)"],
                    size=11)

    # Right: funnel table
    add_text_box(slide, 6.5, 1.2, 6.5, 0.35,
                 "Funil de Triagem por Gene", size=14, bold=True, color=NAVY)

    h_funnel = ["Gene", "Candidatos", "Pass básico", "%"]
    rows_funnel = [
        ("nuc",   "529",    "125",   "23.6%"),
        ("rmpM",  "76",     "32",    "42.1%"),
        ("lytA",  "925",    "540",   "58.4%"),
        ("oprL",  "1 015",  "295",   "29.1%"),
        ("algD",  "2",      "0",     "0.0%"),
        ("frdB",  "1 716",  "467",   "27.2%"),
        ("TOTAL", "4 263",  "1 459", "34.2%"),
    ]
    total_bg  = RGBColor(0xBD, 0xD7, 0xEE)
    fc = [None, None, None, None, ROW_LOW, None, total_bg]
    make_table(slide, 6.5, 1.6, 6.5, 3.3,
               h_funnel, rows_funnel,
               col_widths=[1.3, 1.7, 1.7, 1.8],
               row_colors=fc,
               header_size=11, row_size=11)

    note_box(slide, 6.5, 5.05, 6.5, 0.6,
             "34.2% das candidatas passam todos os critérios termoquímicos. "
             "algD: 0 probes (ver Slide 4).",
             size=10)


# ---------------------------------------------------------------------------
# SLIDE 6 — PASSO 5: ESTRUTURA SECUNDÁRIA (SEQFOLD)
# ---------------------------------------------------------------------------
def slide_seqfold(prs):
    slide = blank_slide(prs)
    set_bg(slide, GREY_BG)
    add_title_bar(slide, "Passo 5: Estrutura Secundária — seqfold",
                  "Previsão de estrutura secundária ssDNA (algoritmo de Nussinov, 37 °C)")

    # Methodology
    add_text_box(slide, 0.3, 1.2, 8.0, 0.35, "Metodologia", size=14, bold=True, color=NAVY)
    lines = [
        "• seqfold: implementação do algoritmo de Nussinov para previsão de\n"
        "  estrutura secundária de ácidos nucleicos",
        "• Temperatura: 37 °C (temperatura fisiológica)",
        "• Calcula: Minimum Free Energy (ΔG MFE) e fração de bases emparelhadas",
        "• Critério de rejeição: ΔG MFE ≤ −6.0 kcal/mol  (Zadeh et al. 2011)\n"
        "  → probe com estrutura secundária estável demais não hibridiza eficientemente",
    ]
    add_bullet_text(slide, 0.3, 1.6, 8.0, 3.0, lines, size=12)

    # Results box
    add_text_box(slide, 0.3, 4.65, 8.0, 0.35, "Resultados", size=14, bold=True, color=NAVY)
    h_sf = ["Gene", "Pass básico", "Pass seqfold", "Rejeitadas", "Threshold"]
    rows_sf = [
        ("nuc",   "125", "118", "7",  "ΔG ≤ −6.0"),
        ("rmpM",  "32",  "32",  "0",  "ΔG ≤ −6.0"),
        ("lytA",  "540", "538", "2",  "ΔG ≤ −6.0"),
        ("oprL",  "295", "295", "0",  "ΔG ≤ −6.0"),
        ("frdB",  "467", "467", "0",  "ΔG ≤ −6.0"),
        ("TOTAL", "1 459", "1 450", "9", "—"),
    ]
    total_bg = RGBColor(0xBD, 0xD7, 0xEE)
    fc = [None, None, None, None, None, total_bg]
    make_table(slide, 0.3, 5.05, 8.0, 2.1,
               h_sf, rows_sf,
               col_widths=[1.2, 1.5, 1.5, 1.4, 2.4],
               row_colors=fc,
               header_size=10, row_size=10)

    # Right: note
    note_box(slide, 8.7, 1.2, 4.3, 2.2,
             "Observação: o threshold −6.0 kcal/mol\n"
             "é demasiado permissivo para ssDNA —\n"
             "filtra apenas 9/1 459 (0.6%) das probes.\n\n"
             "Recomendação: recalibrar para\n"
             "−3.0 kcal/mol em iteração futura.",
             size=11)

    add_text_box(slide, 8.7, 3.55, 4.3, 0.35,
                 "Distribuição de ΔG MFE", size=13, bold=True, color=NAVY)
    add_bullet_text(slide, 8.7, 3.9, 4.3, 2.5,
                    ["• Maioria das probes: ΔG entre 0 e −4 kcal/mol",
                     "• Probes nuc (AT-rich): ΔG tipicamente próximo de 0",
                     "• Probes oprL (GC-rich): ΔG mais negativo (até −5)",
                     "• Nenhuma probe com ΔG < −6 kcal/mol encontrada\n"
                     "  em rmpM, oprL ou frdB"],
                    size=11)


# ---------------------------------------------------------------------------
# SLIDE 7 — PASSO 6: BOLTZ-2 METODOLOGIA
# ---------------------------------------------------------------------------
def slide_boltz2_method(prs):
    slide = blank_slide(prs)
    set_bg(slide, GREY_BG)
    add_title_bar(slide, "Passo 6: Boltz-2 — Previsão de Estrutura 3D",
                  "Rede neuronal de difusão multimodal para ssDNA (Wohlwend et al. 2024)")

    # Left column: methodology
    add_text_box(slide, 0.3, 1.2, 6.2, 0.35, "Modelo Boltz-2", size=14, bold=True, color=NAVY)
    lines_l = [
        "• Boltz-2: rede neuronal de difusão multimodal treinada para\n"
        "  previsão de estrutura 3D de biomoléculas (proteínas e ácidos nucleicos)",
        "• Input: ficheiro YAML com sequência ssDNA\n"
        '  {"version": 1, "sequences": [{"dna": {"id": ["A"], "sequence": "..."}}]}',
        "• Executado em Google Colab com GPU T4\n"
        "  (~2–5 min/probe vs. >30 min em CPU local)",
        "• Adaptação necessária: Boltz-2 requer modo DNA-only\n"
        "  (notebook personalizado — descarta validação de proteína)",
    ]
    add_bullet_text(slide, 0.3, 1.6, 6.2, 3.1, lines_l, size=11)

    # Metrics box
    add_text_box(slide, 0.3, 4.75, 6.2, 0.35,
                 "Métricas de Qualidade Estrutural", size=14, bold=True, color=NAVY)
    h_met = ["Métrica", "Descrição", "Classificação"]
    rows_met = [
        ("confidence", "Score global de confiança da estrutura prevista", "≥ 0.70 → GOOD"),
        ("pTM",        "Predicted TM-score (similaridade estrutural global)", "informativo"),
        ("pLDDT",      "Per-residue Local Distance Difference Test (média)", "≥ 0.80 → alta qualidade"),
    ]
    make_table(slide, 0.3, 5.15, 6.2, 1.85,
               h_met, rows_met,
               col_widths=[1.5, 3.3, 1.4],
               header_size=10, row_size=10)

    # Right column: selection and tiers
    add_text_box(slide, 6.9, 1.2, 6.1, 0.35,
                 "Seleção das 25 Probes", size=14, bold=True, color=NAVY)
    lines_r = [
        "• Probes selecionadas: top 5 por gene (pass básico + melhor seqfold ΔG)",
        "• algD excluído (0 probes disponíveis)",
        "• Total: 5 × 5 genes = 25 probes submetidas ao Boltz-2",
    ]
    add_bullet_text(slide, 6.9, 1.6, 6.1, 1.5, lines_r, size=12)

    # Quality tiers boxes
    add_text_box(slide, 6.9, 3.2, 6.1, 0.35,
                 "Classificação por Confidence Score", size=14, bold=True, color=NAVY)

    tier_data = [
        ("GOOD",     "≥ 0.70",        "2 probes",  ROW_GOOD),
        ("MODERATE", "0.60 – 0.70",   "8 probes",  ROW_MOD),
        ("LOW",      "< 0.60",        "15 probes", ROW_LOW),
    ]
    ty = 3.65
    for tier, rng, count, bg in tier_data:
        shp = add_rect(slide, 6.9, ty, 6.1, 0.65, bg=bg, line_color=MID_GREY)
        tf  = shp.text_frame
        tf.margin_left  = I(0.15)
        tf.margin_top   = I(0.08)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{tier}  ({rng})  →  {count}"
        r.font.size  = P(13)
        r.font.bold  = True if tier == "GOOD" else False
        r.font.color.rgb = DARK_GREY
        ty += 0.72

    note_box(slide, 6.9, 5.85, 6.1, 0.7,
             "Nota: confidence scores para ssDNA tendem a ser menores que para\n"
             "proteínas (Boltz-2 foi treinado maioritariamente em dados proteicos).",
             size=10)


# ---------------------------------------------------------------------------
# SLIDE 8 — RESULTADOS BOLTZ-2 (TABELA COMPLETA)
# ---------------------------------------------------------------------------
BOLTZ_DATA = [
    ("p1065_lytA",  "lytA",  "66.5", "57", "0.636", "0.294", "0.722", "MODERATE"),
    ("p1073_lytA",  "lytA",  "66.8", "62", "0.496", "0.260", "0.555", "LOW"),
    ("p1074_lytA",  "lytA",  "67.6", "59", "0.572", "0.268", "0.649", "LOW"),
    ("p1075_lytA",  "lytA",  "67.8", "61", "0.629", "0.241", "0.726", "MODERATE"),
    ("p1382_lytA",  "lytA",  "67.1", "57", "0.705", "0.319", "0.802", "GOOD"),
    ("p1877_oprL",  "oprL",  "70.6", "68", "0.616", "0.330", "0.687", "MODERATE"),
    ("p205_nuc",    "nuc",   "60.3", "43", "0.430", "0.340", "0.452", "LOW"),
    ("p214_nuc",    "nuc",   "59.8", "46", "0.525", "0.250", "0.593", "LOW"),
    ("p215_nuc",    "nuc",   "60.1", "44", "0.511", "0.263", "0.572", "LOW"),
    ("p216_nuc",    "nuc",   "60.3", "43", "0.495", "0.303", "0.543", "LOW"),
    ("p2194_oprL",  "oprL",  "70.6", "69", "0.581", "0.286", "0.654", "LOW"),
    ("p2195_oprL",  "oprL",  "71.3", "67", "0.532", "0.275", "0.597", "LOW"),
    ("p2196_oprL",  "oprL",  "71.5", "68", "0.512", "0.285", "0.569", "LOW"),
    ("p2205_oprL",  "oprL",  "70.8", "69", "0.556", "0.247", "0.634", "LOW"),
    ("p3367_frdB",  "frdB",  "65.2", "54", "0.792", "0.374", "0.897", "GOOD"),
    ("p3597_frdB",  "frdB",  "64.5", "52", "0.512", "0.292", "0.567", "LOW"),
    ("p3598_frdB",  "frdB",  "64.8", "54", "0.524", "0.299", "0.581", "LOW"),
    ("p3630_frdB",  "frdB",  "64.9", "52", "0.387", "0.217", "0.429", "LOW"),
    ("p3631_frdB",  "frdB",  "64.9", "50", "0.387", "0.280", "0.414", "LOW"),
    ("p498_nuc",    "nuc",   "59.7", "43", "0.626", "0.296", "0.708", "MODERATE"),
    ("p570_rmpM",   "rmpM",  "61.8", "46", "0.535", "0.232", "0.611", "LOW"),
    ("p578_rmpM",   "rmpM",  "60.0", "48", "0.639", "0.259", "0.734", "MODERATE"),
    ("p579_rmpM",   "rmpM",  "60.2", "46", "0.672", "0.343", "0.754", "MODERATE"),
    ("p580_rmpM",   "rmpM",  "60.4", "44", "0.609", "0.284", "0.690", "MODERATE"),
    ("p581_rmpM",   "rmpM",  "61.1", "46", "0.635", "0.352", "0.705", "MODERATE"),
]

QUALITY_COLOR = {"GOOD": ROW_GOOD, "MODERATE": ROW_MOD, "LOW": ROW_LOW}
GOOD_STAR = {"p1382_lytA": " ★", "p3367_frdB": " ★★"}


def slide_boltz2_results(prs):
    slide = blank_slide(prs)
    set_bg(slide, GREY_BG)
    add_title_bar(slide, "Resultados Boltz-2 — 25 Probes",
                  "Previsão de estrutura 3D  |  GOOD ≥ 0.70  ·  MODERATE 0.60–0.70  ·  LOW < 0.60")

    headers = ["Probe ID", "Gene", "Tm (°C)", "GC%", "Confidence", "pTM", "pLDDT", "Qualidade"]
    col_widths = [2.9, 0.75, 0.85, 0.65, 1.15, 0.8, 0.8, 1.2]

    rows = []
    row_colors = []
    for r in BOLTZ_DATA:
        probe_id, gene, tm, gc, conf, ptm, plddt, qual = r
        star = GOOD_STAR.get(probe_id, "")
        rows.append((probe_id, gene, tm, gc, conf, ptm, plddt, qual + star))
        row_colors.append(QUALITY_COLOR.get(qual, None))

    make_table(slide, 0.2, 1.2, 12.9, 6.05,
               headers, rows,
               col_widths=col_widths,
               row_colors=row_colors,
               header_size=10, row_size=8.5)


# ---------------------------------------------------------------------------
# SLIDE 9 — DISCUSSÃO E PRÓXIMOS PASSOS
# ---------------------------------------------------------------------------
def slide_discussion(prs):
    slide = blank_slide(prs)
    set_bg(slide, GREY_BG)
    add_title_bar(slide, "Discussão e Próximos Passos",
                  "Principais conclusões da análise computacional")

    # Left: key results
    add_text_box(slide, 0.3, 1.2, 6.2, 0.35, "Principais Resultados", size=14, bold=True, color=NAVY)

    results = [
        "Pipeline gerou 4 263 candidatas (5 genes) → 1 450 passam triagem completa",
        "",
        "2 Probes GOOD identificadas pelo Boltz-2:",
        "  ★★ p3367_frdB  (H. influenzae)",
        "       confidence=0.792  ·  pLDDT=0.897  →  estrutura local bem definida",
        "  ★   p1382_lytA  (S. pneumoniae)",
        "       confidence=0.705  ·  pLDDT=0.802  ·  Tm=67.1°C  ·  GC=57%",
        "",
        "Gene nuc (S. aureus): todas as 5 probes em LOW",
        "  → AT-rich (GC 43–46%) reduz confidence Boltz-2, MAS favorece hibridação",
        "  → comportamento esperado e desejável para probe de deteção",
        "",
        "Gene algD: problema crítico — apenas 2 seq. no NCBI → 0 probes",
        "  → necessário ampliar query (termos alternativos, mais organismos)",
        "",
        "seqfold filtra apenas 0.6% das probes (9/1 459)",
        "  → threshold −6.0 kcal/mol inadequado para ssDNA curto (18–28 nt)",
    ]
    add_bullet_text(slide, 0.3, 1.6, 6.2, 5.5, results, size=10.5)

    # Right: next steps
    add_text_box(slide, 6.9, 1.2, 6.1, 0.35, "Próximos Passos", size=14, bold=True, color=NAVY)

    steps = [
        ("Especificidade (BLAST)",
         "Verificar off-target binding das 2 probes GOOD contra\ngenoma humano e microbioma oral",
         RGBColor(0xDD, 0xEA, 0xF5)),
        ("Recalibrar seqfold",
         "Testar threshold −3.0 kcal/mol para eliminar probes\ncom estrutura secundária mais estável",
         ROW_MOD),
        ("Resolver algD",
         "Ampliar query NCBI (incluir P. fluorescens, P. putida)\n"
         "ou usar base de dados alternativa (UniProt, SILVA)",
         ROW_LOW),
    ]

    sy = 1.6
    for step_title, step_desc, bg in steps:
        shp = add_rect(slide, 6.9, sy, 6.1, 1.35, bg=bg, line_color=MID_GREY)
        tf  = shp.text_frame
        tf.word_wrap = True
        tf.margin_left  = I(0.12)
        tf.margin_top   = I(0.1)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = step_title
        r.font.size  = P(12)
        r.font.bold  = True
        r.font.color.rgb = DARK_GREY
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = step_desc
        r2.font.size  = P(10)
        r2.font.color.rgb = RGBColor(0x40, 0x40, 0x40)
        sy += 1.45

    # Summary banner
    bar = add_rect(slide, 0, 6.7, W, 0.65, bg=NAVY)
    add_text_box(slide, 0, 6.73, W, 0.55,
                 "Candidatas prioritárias para validação:  p3367_frdB (frdB, H. influenzae)  e  "
                 "p1382_lytA (lytA, S. pneumoniae)",
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# MAIN
# ---------------------------------------------------------------------------
def main():
    out_path = os.path.join("docs", "Apresentacao_Pipeline_GFET.pptx")
    os.makedirs("docs", exist_ok=True)

    prs = new_prs()

    slide_title(prs)
    slide_targets(prs)
    slide_pipeline_overview(prs)
    slide_steps123(prs)
    slide_thermo(prs)
    slide_seqfold(prs)
    slide_boltz2_method(prs)
    slide_boltz2_results(prs)
    slide_discussion(prs)

    prs.save(out_path)
    print(f"Saved: {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
